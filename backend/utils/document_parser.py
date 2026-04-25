"""
====================================================
多格式文档解析器
====================================================
把各种格式的文件提取成纯文本。
支持的格式：PDF / DOCX / DOC / TXT / MD / PPTX

工作原理：
    根据文件扩展名选择对应的解析方法
    PDF  → PyMuPDF（速度快，质量好）
    DOCX → python-docx（直接读 Word 的 XML 结构）
    DOC  → 尝试 PyMuPDF 兜底（旧格式支持有限）
    TXT/MD → 直接读文件内容（自动检测编码）
    PPTX → python-pptx（读取每页的文本框）
"""

import io
import hashlib
from pathlib import Path


class DocumentParser:
    """统一文档解析器"""

    # 扩展名 → 解析方法的映射表
    # 例如：".pdf" → "_parse_pdf"
    SUPPORTED_TYPES = {
        ".pdf": "_parse_pdf",
        ".docx": "_parse_docx",
        ".doc": "_parse_doc",
        ".txt": "_parse_txt",
        ".md": "_parse_txt",
        ".pptx": "_parse_pptx",
    }

    @classmethod
    def supported_extensions(cls) -> list[str]:
        """获取支持的所有文件类型（返回扩展名列表）"""
        return list(cls.SUPPORTED_TYPES.keys())

    @classmethod
    def parse(cls, file: io.IOBase, filename: str) -> str:
        """
        解析文件，返回纯文本

        参数：
            file: 文件流（BytesIO 或文件对象，指针在开头）
            filename: 文件名（用于判断文件类型）

        返回：文件中的纯文本内容
        """
        ext = Path(filename).suffix.lower()  # 取扩展名并转小写，如 ".PDF" → ".pdf"
        parser = cls.SUPPORTED_TYPES.get(ext)
        if not parser:
            raise ValueError(
                f"不支持的文件类型: {ext}，"
                f"支持: {cls.supported_extensions()}"
            )
        return getattr(cls, parser)(file)  # 根据扩展名调用对应的解析方法

    @classmethod
    def _parse_pdf(cls, file):
        """
        解析 PDF 文件
        使用 PyMuPDF（fitz）库，逐页提取文字
        """
        import fitz  # PyMuPDF 库
        doc = fitz.open(stream=file.read(), filetype="pdf")
        pages = [page.get_text() for page in doc]
        doc.close()
        return "\n\n".join(pages)  # 每页之间用空行分隔

    @classmethod
    def _parse_docx(cls, file):
        """
        解析 DOCX 文件
        使用 python-docx 库，读取所有段落文本
        """
        import docx
        doc = docx.Document(file)
        return "\n\n".join([
            p.text for p in doc.paragraphs if p.text.strip()
        ])

    @classmethod
    def _parse_doc(cls, file):
        """
        解析旧版 DOC 文件
        旧版 .doc 是二进制格式，python-docx 不支持
        尝试用 PyMuPDF 当 PDF 解析（有时能成功）
        """
        try:
            return cls._parse_pdf(file)
        except Exception:
            raise ValueError("无法解析 .doc 格式，建议转换为 .docx")

    @classmethod
    def _parse_txt(cls, file):
        """
        解析纯文本文件（TXT / MD）
        需要自动检测编码：utf-8 → gbk → gb2312 → latin-1
        （中文文件常见编码）
        """
        raw = file.read()
        for enc in ["utf-8", "gbk", "gb2312", "latin-1"]:
            try:
                return raw.decode(enc)
            except UnicodeDecodeError:
                continue
        raise ValueError("无法解码文本文件，尝试了 utf-8/gbk/gb2312/latin-1 均失败")

    @classmethod
    def _parse_pptx(cls, file):
        """
        解析 PPTX 文件
        使用 python-pptx 库，读取每页幻灯片中的文本框内容
        """
        from pptx import Presentation
        prs = Presentation(file)
        slides = []
        for slide in prs.slides:
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.paragraphs:
                    for p in shape.text_frame.paragraphs:
                        if p.text.strip():
                            texts.append(p.text)
            if texts:
                slides.append("\n".join(texts))
        return "\n\n---\n\n".join(slides)  # 幻灯片之间用分隔线

    @classmethod
    def compute_hash(cls, file: io.IOBase) -> str:
        """
        计算文件内容的 SHA256 哈希值

        用途：文件去重
        相同内容的文件 → 相同的哈希值 → 不允许重复上传

        参数：
            file: 文件流

        返回：64 位十六进制字符串，如 "a1b2c3d4..."
        """
        file.seek(0)  # 确保从文件开头读
        h = hashlib.sha256()
        # 分块读取（避免大文件一次性占满内存）
        for chunk in iter(lambda: file.read(8192), b""):
            h.update(chunk)
        file.seek(0)  # 把指针放回开头（方便后续再读取）
        return h.hexdigest()
